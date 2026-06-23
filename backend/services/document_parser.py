"""
Multi-format document content extraction service.

Provides parsers for PDF, DOCX, XLSX, XLS, CSV, PPTX, TXT, RTF, ODT,
and image files.  Each parser returns a consistent dict with keys:
text, tables, metadata, pages.
"""

from __future__ import annotations

import logging
from pathlib import Path
from typing import Any

from models.schemas import DocumentCategory

logger = logging.getLogger(__name__)


def _safe_str(value: Any) -> str:
    """Convert a cell value to string, handling None."""
    if value is None:
        return ""
    return str(value)


class DocumentParser:
    """Extracts text and tables from various document formats."""

    def parse(self, file_path: Path, category: DocumentCategory) -> dict[str, Any]:
        """Route to the appropriate parser based on *category* and extension.

        Returns:
            dict with keys: text (str), tables (list[dict]), metadata (dict), pages (int).
        """
        ext = file_path.suffix.lower()
        parsers = {
            ".pdf": self._parse_pdf,
            ".docx": self._parse_docx,
            ".doc": self._parse_docx,  # will fail gracefully
            ".xlsx": self._parse_xlsx,
            ".xls": self._parse_xls,
            ".csv": self._parse_csv,
            ".pptx": self._parse_pptx,
            ".ppt": self._parse_pptx,
            ".txt": self._parse_txt,
            ".rtf": self._parse_rtf,
            ".odt": self._parse_odt,
        }

        if category == DocumentCategory.IMAGE:
            return self._parse_image(file_path)

        parser_fn = parsers.get(ext)
        if parser_fn is None:
            logger.warning("No parser for extension '%s', falling back to TXT.", ext)
            parser_fn = self._parse_txt

        try:
            return parser_fn(file_path)
        except Exception as exc:
            logger.error("Parser failed for '%s': %s", file_path.name, exc)
            return {"text": "", "tables": [], "metadata": {"error": str(exc)}, "pages": 0}

    # ------------------------------------------------------------------
    # Individual parsers
    # ------------------------------------------------------------------

    def _parse_pdf(self, file_path: Path) -> dict[str, Any]:
        """Extract text per page and tables from PDF using PyMuPDF + pdfplumber."""
        import fitz  # PyMuPDF

        doc = fitz.open(str(file_path))
        page_texts: list[str] = []
        for page in doc:
            page_texts.append(page.get_text())

        metadata = dict(doc.metadata) if doc.metadata else {}
        page_count = len(doc)
        doc.close()

        # Try table extraction with pdfplumber
        tables: list[dict[str, Any]] = []
        try:
            import pdfplumber

            with pdfplumber.open(str(file_path)) as pdf:
                for i, page in enumerate(pdf.pages):
                    for tbl in page.extract_tables() or []:
                        if tbl and len(tbl) >= 2:
                            headers = [_safe_str(h) for h in tbl[0]]
                            rows = [[_safe_str(c) for c in row] for row in tbl[1:]]
                            tables.append({"headers": headers, "rows": rows, "page": i + 1})
        except Exception as exc:
            logger.debug("pdfplumber table extraction failed: %s", exc)

        full_text = "\n\n".join(page_texts)
        return {"text": full_text, "tables": tables, "metadata": metadata, "pages": page_count}

    def _parse_docx(self, file_path: Path) -> dict[str, Any]:
        """Extract paragraphs and tables from DOCX."""
        from docx import Document

        doc = Document(str(file_path))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        full_text = "\n".join(paragraphs)

        tables: list[dict[str, Any]] = []
        for table in doc.tables:
            rows_data: list[list[str]] = []
            for row in table.rows:
                rows_data.append([cell.text.strip() for cell in row.cells])
            if rows_data:
                headers = rows_data[0]
                rows = rows_data[1:] if len(rows_data) > 1 else []
                tables.append({"headers": headers, "rows": rows})

        metadata: dict[str, Any] = {}
        if doc.core_properties:
            cp = doc.core_properties
            if cp.author:
                metadata["author"] = cp.author
            if cp.title:
                metadata["title"] = cp.title
            if cp.created:
                metadata["created"] = str(cp.created)

        return {"text": full_text, "tables": tables, "metadata": metadata, "pages": 1}

    def _parse_xlsx(self, file_path: Path) -> dict[str, Any]:
        """Extract all sheets as tables from XLSX."""
        from openpyxl import load_workbook

        wb = load_workbook(str(file_path), read_only=True, data_only=True)
        tables: list[dict[str, Any]] = []
        all_text: list[str] = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows_data: list[list[str]] = []
            for row in ws.iter_rows(values_only=True):
                row_strs = [_safe_str(c) for c in row]
                if any(v.strip() for v in row_strs):
                    rows_data.append(row_strs)

            if rows_data:
                headers = rows_data[0]
                rows = rows_data[1:] if len(rows_data) > 1 else []
                tables.append({"title": sheet_name, "headers": headers, "rows": rows})
                for r in rows_data:
                    all_text.append("\t".join(r))

        wb.close()
        return {
            "text": "\n".join(all_text),
            "tables": tables,
            "metadata": {"sheets": wb.sheetnames},
            "pages": len(wb.sheetnames),
        }

    def _parse_xls(self, file_path: Path) -> dict[str, Any]:
        """Extract data from legacy XLS files."""
        import xlrd

        wb = xlrd.open_workbook(str(file_path))
        tables: list[dict[str, Any]] = []
        all_text: list[str] = []

        for sheet in wb.sheets():
            rows_data: list[list[str]] = []
            for row_idx in range(sheet.nrows):
                row_strs = [_safe_str(sheet.cell_value(row_idx, col)) for col in range(sheet.ncols)]
                rows_data.append(row_strs)

            if rows_data:
                headers = rows_data[0]
                rows = rows_data[1:] if len(rows_data) > 1 else []
                tables.append({"title": sheet.name, "headers": headers, "rows": rows})
                for r in rows_data:
                    all_text.append("\t".join(r))

        return {
            "text": "\n".join(all_text),
            "tables": tables,
            "metadata": {"sheets": [s.name for s in wb.sheets()]},
            "pages": len(wb.sheets()),
        }

    def _parse_csv(self, file_path: Path) -> dict[str, Any]:
        """Read CSV into a table using pandas with smart detection."""
        import pandas as pd

        # Read raw content first to detect format
        raw_bytes = file_path.read_bytes()
        for enc in ("utf-8", "utf-8-sig", "latin-1", "cp1252"):
            try:
                raw_text = raw_bytes.decode(enc)
                break
            except UnicodeDecodeError:
                continue
        else:
            raw_text = raw_bytes.decode("utf-8", errors="replace")

        # Auto-detect delimiter
        delimiter = ","
        first_lines = raw_text.strip().split("\n")[:5]
        sample = "\n".join(first_lines)
        if sample.count("\t") > sample.count(","):
            delimiter = "\t"
        elif sample.count(";") > sample.count(","):
            delimiter = ";"

        try:
            df = pd.read_csv(str(file_path), dtype=str, keep_default_na=False,
                             sep=delimiter, encoding=enc)
        except Exception:
            # Fallback: read without header assumption
            df = pd.read_csv(str(file_path), dtype=str, keep_default_na=False,
                             sep=delimiter, header=None, encoding=enc)

        # If only 1 row or 0 rows, the first row might actually be data
        # Re-read with header=None and use generic column names
        if len(df) == 0 and len(df.columns) >= 2:
            df = pd.read_csv(str(file_path), dtype=str, keep_default_na=False,
                             sep=delimiter, header=None, encoding=enc)
            df.columns = [f"Column_{i+1}" for i in range(len(df.columns))]

        headers = [str(h) for h in df.columns]
        rows = df.values.tolist()
        display_text = df.to_string(index=False)

        return {
            "text": display_text,
            "tables": [{"title": file_path.stem, "headers": headers, "rows": rows}],
            "metadata": {"rows": len(df), "columns": len(headers), "delimiter": delimiter},
            "pages": 1,
        }

    def _parse_pptx(self, file_path: Path) -> dict[str, Any]:
        """Extract text and tables from PPTX presentations."""
        from pptx import Presentation

        prs = Presentation(str(file_path))
        slide_texts: list[str] = []
        tables: list[dict[str, Any]] = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            texts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            texts.append(para.text.strip())

                if shape.has_table:
                    tbl = shape.table
                    rows_data: list[list[str]] = []
                    for row in tbl.rows:
                        rows_data.append([cell.text.strip() for cell in row.cells])
                    if rows_data:
                        headers = rows_data[0]
                        rows = rows_data[1:] if len(rows_data) > 1 else []
                        tables.append({"title": f"Slide {slide_num}", "headers": headers, "rows": rows, "page": slide_num})

            slide_texts.append(f"--- Slide {slide_num} ---\n" + "\n".join(texts))

        return {
            "text": "\n\n".join(slide_texts),
            "tables": tables,
            "metadata": {"slide_count": len(prs.slides)},
            "pages": len(prs.slides),
        }

    def _parse_txt(self, file_path: Path) -> dict[str, Any]:
        """Read plain text with multi-encoding fallback."""
        text = ""
        for encoding in ("utf-8", "latin-1", "cp1252"):
            try:
                text = file_path.read_text(encoding=encoding)
                break
            except (UnicodeDecodeError, ValueError):
                continue

        return {"text": text, "tables": [], "metadata": {"encoding": encoding}, "pages": 1}

    def _parse_rtf(self, file_path: Path) -> dict[str, Any]:
        """Convert RTF to plain text."""
        from striprtf.striprtf import rtf_to_text

        raw = file_path.read_bytes().decode("utf-8", errors="replace")
        text = rtf_to_text(raw)
        return {"text": text, "tables": [], "metadata": {}, "pages": 1}

    def _parse_odt(self, file_path: Path) -> dict[str, Any]:
        """Extract text from OpenDocument Text files."""
        from odf.opendocument import load
        from odf import text as odf_text
        from odf.text import P

        doc = load(str(file_path))
        paragraphs: list[str] = []
        for p in doc.getElementsByType(P):
            content = ""
            for node in p.childNodes:
                if hasattr(node, "data"):
                    content += node.data
                elif hasattr(node, "__str__"):
                    content += str(node)
            if content.strip():
                paragraphs.append(content.strip())

        return {"text": "\n".join(paragraphs), "tables": [], "metadata": {}, "pages": 1}

    def _parse_image(self, file_path: Path) -> dict[str, Any]:
        """Return image metadata — actual text extraction happens via OCR/AI."""
        try:
            from PIL import Image

            img = Image.open(str(file_path))
            metadata = {
                "width": img.width,
                "height": img.height,
                "format": img.format,
                "mode": img.mode,
            }
            img.close()
        except Exception:
            metadata = {}

        return {"text": "", "tables": [], "metadata": metadata, "pages": 1}
